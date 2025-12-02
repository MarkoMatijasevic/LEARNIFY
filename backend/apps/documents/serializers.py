from rest_framework import serializers
from .models import Document, DocumentCategory, DocumentShare, DocumentTest, TestAttempt
import os
import logging
import PyPDF2
import docx
from pptx import Presentation
import io

# Set up logging
logger = logging.getLogger(__name__)

class DocumentCategorySerializer(serializers.ModelSerializer):
    class Meta:
        model = DocumentCategory
        fields = ['id', 'name', 'description', 'color', 'icon']

class DocumentSerializer(serializers.ModelSerializer):
    category_name = serializers.CharField(source='category.name', read_only=True)
    file_size_mb = serializers.SerializerMethodField()
    is_ready = serializers.SerializerMethodField()
    tags_list = serializers.SerializerMethodField()
    
    class Meta:
        model = Document
        fields = [
            'id', 'title', 'description', 'original_filename', 
            'file_size', 'file_size_mb', 'file_type', 'status',
            'category', 'category_name', 'tags', 'tags_list',
            'text_preview', 'page_count', 'word_count',
            'ai_summary', 'ai_key_topics', 'ai_difficulty_level',
            'created_at', 'updated_at', 'last_accessed', 'is_ready'
        ]
        read_only_fields = [
            'id', 'file_size', 'file_type', 'status', 'extracted_text',
            'text_preview', 'page_count', 'word_count', 'ai_summary',
            'ai_key_topics', 'ai_difficulty_level', 'created_at', 
            'updated_at', 'last_accessed'
        ]
    
    def get_file_size_mb(self, obj):
        return obj.file_size_mb
    
    def get_is_ready(self, obj):
        return obj.is_ready
    
    def get_tags_list(self, obj):
        return obj.tags_list

class DocumentUploadSerializer(serializers.ModelSerializer):
    file = serializers.FileField()
    category_id = serializers.IntegerField(required=False, allow_null=True)
    
    class Meta:
        model = Document
        fields = ['title', 'description', 'file', 'category_id', 'tags']
    
    def validate_file(self, value):
        # Check file size (50MB limit)
        if value.size > 50 * 1024 * 1024:
            raise serializers.ValidationError("File size cannot exceed 50MB")
        
        # Check file type
        allowed_extensions = ['.pdf', '.docx', '.doc', '.pptx', '.ppt', '.txt']
        file_extension = os.path.splitext(value.name)[1].lower()
        
        if file_extension not in allowed_extensions:
            raise serializers.ValidationError(
                f"File type {file_extension} not supported. "
                f"Allowed types: {', '.join(allowed_extensions)}"
            )
        
        return value
    
    def extract_text_from_pdf(self, file_content):
        """Extract text from PDF file"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text = ""
            page_count = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text += f"\n--- Page {page_num + 1} ---\n{page_text}\n"
                except Exception as e:
                    logger.warning(f"Failed to extract text from page {page_num + 1}: {e}")
                    continue
            
            return text.strip(), page_count
        except Exception as e:
            logger.error(f"PDF text extraction error: {e}")
            return "", 0
    
    def extract_text_from_docx(self, file_content):
        """Extract text from Word document"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = docx.Document(doc_file)
            
            text_parts = []
            
            # Extract paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    text_parts.append(para.text.strip())
            
            # Extract tables
            for table in doc.tables:
                for row in table.rows:
                    row_text = []
                    for cell in row.cells:
                        if cell.text.strip():
                            row_text.append(cell.text.strip())
                    if row_text:
                        text_parts.append(" | ".join(row_text))
            
            text = "\n".join(text_parts)
            
            # Count pages (approximation: 500 words per page)
            word_count = len(text.split())
            page_count = max(1, word_count // 500)
            
            return text, page_count
        except Exception as e:
            logger.error(f"Word document text extraction error: {e}")
            return "", 1
    
    def extract_text_from_pptx(self, file_content):
        """Extract text from PowerPoint presentation"""
        try:
            ppt_file = io.BytesIO(file_content)
            presentation = Presentation(ppt_file)
            
            text_parts = []
            slide_count = 0
            
            for slide_num, slide in enumerate(presentation.slides):
                slide_count += 1
                slide_text = []
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    text_parts.append(f"\n--- Slide {slide_num + 1} ---\n" + "\n".join(slide_text))
            
            text = "\n".join(text_parts)
            return text, slide_count
        except Exception as e:
            logger.error(f"PowerPoint text extraction error: {e}")
            return "", 1
    
    def extract_text_content(self, file, file_type):
        """Extract text content based on file type"""
        try:
            # Read file content
            file.seek(0)
            file_content = file.read()
            file.seek(0)  # Reset file pointer
            
            logger.info(f"Extracting text from {file_type} file, size: {len(file_content)} bytes")
            
            if file_type == 'pdf':
                extracted_text, page_count = self.extract_text_from_pdf(file_content)
            elif file_type in ['docx', 'doc']:
                extracted_text, page_count = self.extract_text_from_docx(file_content)
            elif file_type in ['pptx', 'ppt']:
                extracted_text, page_count = self.extract_text_from_pptx(file_content)
            elif file_type == 'txt':
                try:
                    extracted_text = file_content.decode('utf-8', errors='ignore')
                except UnicodeDecodeError:
                    extracted_text = file_content.decode('latin-1', errors='ignore')
                
                word_count = len(extracted_text.split())
                page_count = max(1, word_count // 500)
            else:
                logger.warning(f"Unsupported file type: {file_type}")
                return "", 0
            
            # Clean and validate extracted text
            if extracted_text:
                extracted_text = extracted_text.strip()
                if len(extracted_text) < 10:  # Minimum content threshold
                    logger.warning(f"Extracted text too short: {len(extracted_text)} characters")
                    return "", page_count
                
                logger.info(f"Successfully extracted {len(extracted_text)} characters from {file_type}")
                return extracted_text, page_count
            else:
                logger.warning(f"No text extracted from {file_type} file")
                return "", page_count
                
        except Exception as e:
            logger.error(f"Text extraction error for {file_type}: {e}", exc_info=True)
            return "", 0
    
    def create(self, validated_data):
        file = validated_data.pop('file')
        category_id = validated_data.pop('category_id', None)
        user = self.context['request'].user
        
        try:
            # Set category if provided
            category = None
            if category_id:
                try:
                    category = DocumentCategory.objects.get(id=category_id)
                except DocumentCategory.DoesNotExist:
                    logger.warning(f"Category with id {category_id} not found")
            
            # Determine file type
            file_extension = os.path.splitext(file.name)[1].lower().replace('.', '')
            file_type_map = {
                'pdf': 'pdf',
                'doc': 'docx', 'docx': 'docx',
                'ppt': 'pptx', 'pptx': 'pptx',
                'txt': 'txt'
            }
            
            file_type = file_type_map.get(file_extension, 'txt')
            
            # Extract text content
            logger.info(f"Processing file: {file.name} (type: {file_type}, size: {file.size} bytes)")
            extracted_text, page_count = self.extract_text_content(file, file_type)
            
            # Calculate word count
            word_count = len(extracted_text.split()) if extracted_text else 0
            
            # Create text preview (first 500 characters)
            text_preview = extracted_text[:500] + "..." if len(extracted_text) > 500 else extracted_text
            
            # Determine status based on extraction success
            status = 'ready' if extracted_text and len(extracted_text) > 10 else 'error'
            
            # Create document instance
            document = Document.objects.create(
                user=user,
                title=validated_data.get('title', os.path.splitext(file.name)[0]),
                description=validated_data.get('description', ''),
                file=file,
                original_filename=file.name,
                file_size=file.size,
                file_type=file_type,
                tags=validated_data.get('tags', ''),
                category=category,
                extracted_text=extracted_text,  # CRITICAL: Ensure this field is set
                text_preview=text_preview,
                page_count=page_count,
                word_count=word_count,
                status=status
            )
            
            # Verify the document was saved with text
            document.refresh_from_db()
            saved_text_length = len(document.extracted_text) if document.extracted_text else 0
            
            logger.info(f"Document {document.id} created successfully:")
            logger.info(f"  - Status: {document.status}")
            logger.info(f"  - Extracted text length: {saved_text_length} characters")
            logger.info(f"  - Word count: {document.word_count}")
            logger.info(f"  - Page count: {document.page_count}")
            
            if saved_text_length == 0:
                logger.error(f"WARNING: Document {document.id} was saved but has no extracted text!")
                # Update status to indicate problem
                document.status = 'error'
                document.save(update_fields=['status'])
            
            return document
            
        except Exception as e:
            logger.error(f"Document creation failed: {str(e)}", exc_info=True)
            raise serializers.ValidationError(f"Document upload failed: {str(e)}")

class DocumentShareSerializer(serializers.ModelSerializer):
    shared_with_email = serializers.CharField(source='shared_with.email', read_only=True)
    document_title = serializers.CharField(source='document.title', read_only=True)
    
    class Meta:
        model = DocumentShare
        fields = [
            'id', 'document', 'document_title', 'shared_with', 
            'shared_with_email', 'permission', 'created_at'
        ]
        read_only_fields = ['id', 'created_at']

# ============================================================================
# TEST GENERATION SERIALIZERS (NEW)
# ============================================================================

class TestQuestionSerializer(serializers.Serializer):
    """Serializer for individual test questions"""
    id = serializers.IntegerField()
    question = serializers.CharField()
    options = serializers.DictField(child=serializers.CharField())
    correct_answer = serializers.ChoiceField(choices=['A', 'B', 'C'])
    explanation = serializers.CharField()


class DocumentTestSerializer(serializers.ModelSerializer):
    """Serializer for DocumentTest model"""
    document_title = serializers.CharField(source='document.title', read_only=True)
    questions = TestQuestionSerializer(many=True, read_only=True)
    is_ready = serializers.SerializerMethodField()
    attempt_count = serializers.SerializerMethodField()
    
    class Meta:
        model = DocumentTest
        fields = [
            'id', 'document', 'document_title', 'created_by',
            'title', 'question_count', 'questions', 'status',
            'generation_error', 'generation_time_seconds',
            'created_at', 'updated_at', 'is_ready', 'attempt_count'
        ]
        read_only_fields = [
            'id', 'created_by', 'questions', 'status',
            'generation_error', 'generation_time_seconds',
            'created_at', 'updated_at'
        ]
    
    def get_is_ready(self, obj):
        return obj.is_ready
    
    def get_attempt_count(self, obj):
        return obj.attempt_count


class TestGenerateRequestSerializer(serializers.Serializer):
    """Serializer for test generation request"""
    document_id = serializers.UUIDField()
    
    def validate_document_id(self, value):
        """Validate that document exists and belongs to user"""
        request = self.context.get('request')
        try:
            document = Document.objects.get(id=value, user=request.user)
            
            # Check if document is ready
            if not document.is_ready:
                raise serializers.ValidationError(
                    "Document is not ready for test generation. "
                    "Please wait for document processing to complete."
                )
            
            # Check if document has enough content
            if not document.extracted_text or len(document.extracted_text.strip()) < 100:
                raise serializers.ValidationError(
                    "Document does not have enough text content for test generation."
                )
            
            return value
            
        except Document.DoesNotExist:
            raise serializers.ValidationError("Document not found or you don't have access to it.")


class TestResultDetailSerializer(serializers.Serializer):
    """Serializer for individual question result details"""
    question_id = serializers.IntegerField()
    question = serializers.CharField()
    options = serializers.DictField(child=serializers.CharField())
    user_answer = serializers.CharField(allow_blank=True)
    correct_answer = serializers.ChoiceField(choices=['A', 'B', 'C'])
    is_correct = serializers.BooleanField()
    explanation = serializers.CharField()


class TestAttemptSerializer(serializers.ModelSerializer):
    """Serializer for TestAttempt model"""
    test_title = serializers.CharField(source='test.title', read_only=True)
    results_detail = TestResultDetailSerializer(many=True, read_only=True)
    score_percentage = serializers.SerializerMethodField()
    time_taken_formatted = serializers.SerializerMethodField()
    
    class Meta:
        model = TestAttempt
        fields = [
            'id', 'test', 'test_title', 'user', 'answers',
            'score', 'score_percentage', 'grade', 'passed',
            'correct_count', 'incorrect_count', 'results_detail',
            'time_taken_seconds', 'time_taken_formatted',
            'started_at', 'completed_at'
        ]
        read_only_fields = [
            'id', 'user', 'score', 'grade', 'passed',
            'correct_count', 'incorrect_count', 'results_detail',
            'started_at', 'completed_at'
        ]
    
    def get_score_percentage(self, obj):
        return obj.score_percentage
    
    def get_time_taken_formatted(self, obj):
        return obj.time_taken_formatted


class TestSubmissionSerializer(serializers.Serializer):
    """Serializer for test submission"""
    answers = serializers.DictField(
        child=serializers.ChoiceField(choices=['A', 'B', 'C', '']),
        help_text="Dictionary mapping question IDs to selected answers (A, B, or C)"
    )
    time_taken_seconds = serializers.IntegerField(
        required=False,
        allow_null=True,
        help_text="Time taken to complete the test in seconds"
    )
    
    def validate_answers(self, value):
        """Validate that answers are in correct format"""
        if not isinstance(value, dict):
            raise serializers.ValidationError("Answers must be a dictionary")
        
        # Validate answer values
        for question_id, answer in value.items():
            if answer and answer not in ['A', 'B', 'C']:
                raise serializers.ValidationError(
                    f"Invalid answer '{answer}' for question {question_id}. "
                    f"Answer must be A, B, or C."
                )
        
        return value
    
    def validate_time_taken_seconds(self, value):
        """Validate time taken is reasonable"""
        if value is not None and value < 0:
            raise serializers.ValidationError("Time taken cannot be negative")
        
        if value is not None and value > 86400:  # 24 hours
            raise serializers.ValidationError("Time taken seems unreasonably long")
        
        return value


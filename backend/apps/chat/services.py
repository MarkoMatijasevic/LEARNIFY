# backend/apps/chat/services.py
import os
import logging
from typing import Dict, List, Optional
import google.generativeai as genai
from django.conf import settings
import PyPDF2
from docx import Document as DocxDocument

logger = logging.getLogger(__name__)

class GeminiService:
    """Service for interacting with Google Gemini AI API"""
    
    def __init__(self, ai_model=None):
        self.ai_model = ai_model
        self.setup_gemini()
    
    def setup_gemini(self):
        """Initialize Gemini AI client"""
        try:
            api_key = getattr(settings, 'GOOGLE_AI_API_KEY', os.getenv('GOOGLE_AI_API_KEY'))
            if not api_key:
                raise ValueError("Google AI API key not found in settings or environment")
            
            genai.configure(api_key=api_key)
            self.model = genai.GenerativeModel('gemini-2.5-flash')
            logger.info("Gemini AI client initialized successfully")
            
        except Exception as e:
            logger.error(f"Failed to initialize Gemini AI: {str(e)}")
            raise
    
    def generate_response(
        self, 
        message: str, 
        document_context: str = "", 
        conversation_history: List[Dict] = None
    ) -> Dict:
        """Generate AI response using Gemini"""
        
        try:
            # Build prompt with context
            system_prompt = self._build_system_prompt()
            context_prompt = self._build_context_prompt(document_context, conversation_history)
            full_prompt = f"{system_prompt}\n\n{context_prompt}\n\nUser: {message}"
            
            # Generate response
            response = self.model.generate_content(full_prompt)
            
            # Calculate approximate token usage (Gemini doesn't provide exact counts)
            tokens_used = self._estimate_tokens(full_prompt + response.text)
            
            return {
                'content': response.text.strip(),
                'tokens_used': tokens_used,
                'success': True
            }
            
        except Exception as e:
            logger.error(f"Gemini API error: {str(e)}")
            return {
                'content': "I'm sorry, I encountered an error while processing your request. Please try again.",
                'tokens_used': 0,
                'success': False,
                'error': str(e)
            }
    
    def _build_system_prompt(self) -> str:
        """Build system prompt for the AI assistant"""
        return """You are Learnify AI, an intelligent learning assistant. Your role is to help users understand and learn from their uploaded study materials.

Guidelines:
- Be helpful, accurate, and educational
- Focus on learning and comprehension
- Use the provided document context to answer questions
- If you don't know something from the documents, say so clearly
- Encourage critical thinking and deeper understanding
- Provide explanations in a clear, structured way
- Use examples when helpful
- Stay on topic and be concise but thorough

Always prioritize accuracy and educational value in your responses."""
    
    def _build_context_prompt(self, document_context: str, conversation_history: List[Dict]) -> str:
        """Build context prompt with documents and history"""
        prompt_parts = []
        
        # Add document context if available
        if document_context:
            prompt_parts.append(f"DOCUMENT CONTEXT:\n{document_context}")
        
        # Add conversation history if available
        if conversation_history:
            history_text = "\n".join([
                f"{msg['role'].title()}: {msg['content']}" 
                for msg in conversation_history[-10:]  # Last 10 messages
            ])
            prompt_parts.append(f"CONVERSATION HISTORY:\n{history_text}")
        
        return "\n\n".join(prompt_parts)
    
    def _estimate_tokens(self, text: str) -> int:
        """Estimate token count (rough approximation)"""
        # Rough estimation: ~4 characters per token
        return len(text) // 4

class DocumentProcessor:
    """Service for extracting text from various document types"""
    
    def extract_text(self, file_path: str) -> str:
        """Extract text from document based on file extension"""
        
        if not os.path.exists(file_path):
            raise FileNotFoundError(f"File not found: {file_path}")
        
        file_extension = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_extension == '.pdf':
                return self._extract_pdf_text(file_path)
            elif file_extension in ['.docx', '.doc']:
                return self._extract_docx_text(file_path)
            elif file_extension in ['.pptx', '.ppt']:
                return self._extract_pptx_text(file_path)
            else:
                raise ValueError(f"Unsupported file type: {file_extension}")
                
        except Exception as e:
            logger.error(f"Error extracting text from {file_path}: {str(e)}")
            raise
    
    def _extract_pdf_text(self, file_path: str) -> str:
        """Extract text from PDF file"""
        text_content = []
        
        try:
            with open(file_path, 'rb') as file:
                pdf_reader = PyPDF2.PdfReader(file)
                
                for page_num, page in enumerate(pdf_reader.pages):
                    try:
                        page_text = page.extract_text()
                        if page_text.strip():
                            text_content.append(f"--- Page {page_num + 1} ---\n{page_text}")
                    except Exception as e:
                        logger.warning(f"Error extracting page {page_num + 1}: {str(e)}")
                        continue
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"PDF extraction error: {str(e)}")
            raise
    
    def _extract_docx_text(self, file_path: str) -> str:
        """Extract text from Word document"""
        try:
            doc = DocxDocument(file_path)
            text_content = []
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content.append(paragraph.text)
            
            return "\n\n".join(text_content)
            
        except Exception as e:
            logger.error(f"DOCX extraction error: {str(e)}")
            raise
    
    def _extract_pptx_text(self, file_path: str) -> str:
        """Extract text from PowerPoint presentation"""
        try:
            from pptx import Presentation
            
            prs = Presentation(file_path)
            text_content = []
            
            for slide_num, slide in enumerate(prs.slides):
                slide_text = [f"--- Slide {slide_num + 1} ---"]
                
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text)
                
                if len(slide_text) > 1:  # Has content beyond slide number
                    text_content.append("\n".join(slide_text))
            
            return "\n\n".join(text_content)
            
        except ImportError:
            logger.error("python-pptx package not installed")
            raise ValueError("PowerPoint processing not available. Install python-pptx package.")
        except Exception as e:
            logger.error(f"PPTX extraction error: {str(e)}")
            raise
    
    def get_document_summary(self, file_path: str, max_chars: int = 500) -> str:
        """Get a summary of the document content"""
        try:
            full_text = self.extract_text(file_path)
            
            if len(full_text) <= max_chars:
                return full_text
            
            # Return first portion of the document
            truncated = full_text[:max_chars]
            last_sentence = truncated.rfind('.')
            
            if last_sentence > max_chars * 0.7:  # If we can find a good sentence break
                return truncated[:last_sentence + 1] + "..."
            else:
                return truncated + "..."
                
        except Exception as e:
            logger.error(f"Error generating document summary: {str(e)}")
            return "Unable to generate document summary."